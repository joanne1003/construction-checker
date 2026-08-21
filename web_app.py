import streamlit as st
import json
import time
from io import BytesIO
import google.generativeai as genai
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor

# 設定 API
genai.configure(api_key=st.secrets["GEMINI_API_KEY"])

@st.cache_resource
def get_working_model():
    try:
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        for pref in ['gemini-1.5-flash', 'gemini-2.0-flash', 'gemini-1.5-pro', 'gemini-pro']:
            for am in available_models:
                if pref in am:
                    return genai.GenerativeModel(am)
        if available_models:
            return genai.GenerativeModel(available_models[0])
    except Exception:
        pass
    
    for name in ['gemini-1.5-flash', 'gemini-1.5-pro', 'gemini-pro']:
        try:
            return genai.GenerativeModel(name)
        except:
            continue
    return genai.GenerativeModel('gemini-1.5-flash')

model = get_working_model()

# 專業工程配色 (完全對齊參考圖風格)
COLOR_BLUE = RGBColor(0, 80, 160)
COLOR_GREEN = RGBColor(0, 128, 64)
COLOR_DARK = RGBColor(40, 40, 40)

st.set_page_config(page_title="專業職安工程檢核圖產生器", page_icon="🏗️")
st.title("🏗️ 專業職安工程檢核圖產生器")

uploaded_file = st.file_uploader("📂 上傳施工照片", type=["jpg", "png"])

if uploaded_file is not None:
    st.image(uploaded_file, use_container_width=True)
    if st.button("🚀 生成與參考圖風格一致的專業檢核圖"):
        with st.spinner('專業職安工程師正在規劃圖說（若遇限流將自動重試）...'):
            try:
                img_pil = Image.open(uploaded_file).convert("RGB")
                
                # 【優化】圖片尺寸自動縮放（限制最大邊長 1200px），大幅降低 Token 消耗、避免觸發 429
                max_dim = 1200
                if max(img_pil.size) > max_dim:
                    img_pil.thumbnail((max_dim, max_dim))
                
                width_px, height_px = img_pil.size
                
                prompt = """
                身為資深職業安全衛生工程師與營造工程專家，請分析此施工照片。
                1. 辨識工程類型 (例如：橋梁墩柱、擋土支撐、瀝青鋪設)。
                2. 找出 4-6 個關鍵的工安風險點或品質檢核點。
                3. 檢核要點必須具體且符合營造安全衛生設施標準。
                
                回傳嚴格的 JSON 格式 (絕對不可包含 markdown 標籤或說明文字)：
                {
                  "project_type": "工程名稱",
                  "checkpoints": [
                    {
                      "label": "1",
                      "box_2d": [ymin, xmin, ymax, xmax],
                      "title": "檢核項目",
                      "items": "要點一\n要點二"
                    }
                  ]
                }
                """
                
                # 【優化】加入自動重試機制 (Auto-retry for 429 Quota Exceeded)
                max_retries = 3
                response = None
                for attempt in range(max_retries):
                    try:
                        response = model.generate_content([img_pil, prompt])
                        break
                    except Exception as api_err:
                        err_str = str(api_err)
                        if ("429" in err_str or "Quota exceeded" in err_str) and attempt < max_retries - 1:
                            time.sleep(15 * (attempt + 1))  # 遞增等待 15秒、30秒
                            continue
                        raise api_err

                text = response.text.replace('```json', '').replace('```', '').strip()
                data = json.loads(text)
                
                proj_type = data.get("project_type", "工程")
                checkpoints = data.get("checkpoints", [])

                # 依據位置分邊，確保畫面整潔不交錯
                left_items = []
                right_items = []
                for item in checkpoints:
                    if (item["box_2d"][1] + item["box_2d"][3]) / 2 < 500:
                        left_items.append(item)
                    else:
                        right_items.append(item)
                
                prs = Presentation()
                prs.slide_height = Inches(7.5) 
                prs.slide_width = Inches(7.5 * (width_px / height_px))
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 背景圖
                img_io = BytesIO()
                img_pil.save(img_io, format='JPEG')
                img_io.seek(0)
                slide.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)

                # 標題區 (圓角矩形)
                t_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.3), Inches(5.2), Inches(0.6))
                t_box.fill.solid()
                t_box.fill.fore_color.rgb = COLOR_BLUE
                tf = t_box.text_frame
                tf.text = f"{proj_type}施工簡易檢核圖"
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                s_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.95), Inches(5.2), Inches(0.45))
                s_box.fill.solid()
                s_box.fill.fore_color.rgb = COLOR_DARK
                tf_sub = s_box.text_frame
                tf_sub.text = f"{proj_type}職業安全衛生與品質檢核要點"
                tf_sub.paragraphs[0].font.size = Pt(13)
                tf_sub.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

                # 繪製卡片
                def draw_cards(items, is_left):
                    base_color = COLOR_BLUE if is_left else COLOR_GREEN
                    card_x = Inches(0.4) if is_left else prs.slide_width - Inches(3.9)
                    
                    for i, item in enumerate(items):
                        card_y = Inches(1.6 + i * 1.6)
                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, Inches(3.5), Inches(1.4))
                        card.fill.solid()
                        card.fill.fore_color.rgb = base_color
                        card.line.color.rgb = RGBColor(255, 255, 255)
                        card.line.width = Pt(1.5)
                        
                        tf = card.text_frame
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                        
                        p = tf.paragraphs[0]
                        p.text = f"{item.get('label', '')}. {item['title']}"
                        p.font.bold = True
                        p.font.size = Pt(13)
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.space_after = Pt(25)

                        line_sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(card_x + Inches(0.15)), int(card_y + Inches(0.6)), int(card_x + Inches(3.35)), int(card_y + Inches(0.6)))
                        line_sep.line.color.rgb = RGBColor(255, 255, 255)
                        line_sep.line.width = Pt(1.5)

                        for it in item["items"].split('\n'):
                            if it.strip():
                                p = tf.add_paragraph()
                                p.text = f"☑ {it.strip()}"
                                p.font.size = Pt(10.5)
                                p.font.color.rgb = RGBColor(255, 255, 255)
                        
                        ymin, xmin, ymax, xmax = item["box_2d"]
                        tx = int((xmin + xmax) / 2 / 1000 * prs.slide_width)
                        ty = int((ymin + ymax) / 2 / 1000 * prs.slide_height)
                        lx = int(card_x + Inches(3.5)) if is_left else int(card_x)
                        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, lx, int(card_y + Inches(0.7)), tx, ty)
                        line.line.color.rgb = RGBColor(255, 255, 255)
                        line.line.width = Pt(2)
                        
                        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, tx-12, ty-12, 24, 24)
                        circle.fill.solid()
                        circle.fill.fore_color.rgb = base_color
                        circle.text_frame.text = item.get('label', '')

                draw_cards(left_items, True)
                draw_cards(right_items, False)

                table = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, prs.slide_width - Inches(3.9), prs.slide_height - Inches(1.8), Inches(3.5), Inches(1.5))
                table.fill.solid()
                table.fill.fore_color.rgb = RGBColor(255, 255, 255)
                table.text_frame.text = "現場檢核結果記錄\n(項目 | 檢核結果 | 備註)"

                pptx_io = BytesIO()
                prs.save(pptx_io)
                pptx_io.seek(0)
                st.success("✅ 專業職安檢核圖生成成功！")
                st.download_button("📥 下載專業檢核報告.pptx", pptx_io, "專業檢核報告.pptx")
            except Exception as e:
                err_msg = str(e)
                if "429" in err_msg or "Quota exceeded" in err_msg:
                    st.warning("⚠️ **已達到免費 API 請求上限 (429 Quota Exceeded)**\n\n系統已嘗試自動重試但仍超過限制。建議**稍候約 1 分鐘**再點擊生成，或是前往 Google AI Studio 檢查您的 API 額度。")
                else:
                    st.error(f"發生錯誤: {err_msg}")
